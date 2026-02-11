from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ankismart.core.errors import ConvertError, ErrorCode
from ankismart.core.logging import get_logger
from ankismart.core.models import MarkdownResult
from ankismart.core.tracing import get_trace_id, timed

logger = get_logger("converter.pptx")


def _get_slide_title(slide) -> str | None:
    if slide.shapes.title is not None:
        return slide.shapes.title.text.strip() or None
    return None


def _render_runs(paragraph) -> str:
    """Render paragraph runs with bold/italic Markdown formatting."""
    parts: list[str] = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        is_bold = run.font.bold
        is_italic = run.font.italic
        if is_bold and is_italic:
            parts.append(f"***{text}***")
        elif is_bold:
            parts.append(f"**{text}**")
        elif is_italic:
            parts.append(f"*{text}*")
        else:
            parts.append(text)
    return "".join(parts) or paragraph.text or ""


def _extract_slide_text(slide) -> list[str]:
    lines: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = _render_runs(paragraph).strip()
            if not text:
                continue
            raw_level = paragraph.level
            level = raw_level if isinstance(raw_level, int) else 0
            if level > 0:
                indent = "  " * level
                lines.append(f"{indent}- {text}")
            else:
                lines.append(text)
    return lines


def convert(file_path: Path, trace_id: str = "") -> MarkdownResult:
    tid = trace_id or get_trace_id()
    with timed("pptx_convert"):
        try:
            prs = Presentation(str(file_path))
        except FileNotFoundError as exc:
            raise ConvertError(
                f"File not found: {file_path}",
                code=ErrorCode.E_FILE_NOT_FOUND,
                trace_id=tid,
            ) from exc
        except Exception as exc:
            raise ConvertError(
                f"Failed to open pptx: {file_path}: {exc}",
                trace_id=tid,
            ) from exc

        slide_parts: list[str] = []

        for idx, slide in enumerate(prs.slides, start=1):
            title = _get_slide_title(slide)
            header = f"## {title}" if title else f"## Slide {idx}"

            lines = _extract_slide_text(slide)
            body = "\n".join(lines)

            slide_parts.append(f"{header}\n\n{body}" if body else header)

        content = "\n\n---\n\n".join(slide_parts) + "\n"
        logger.info("Converted pptx file", extra={"path": str(file_path), "trace_id": tid})

        return MarkdownResult(
            content=content,
            source_path=str(file_path),
            source_format="pptx",
            trace_id=tid,
        )
